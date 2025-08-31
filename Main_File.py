import pptx
import re
import io
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image, ImageEnhance


# Function to color the word "Google". This was for all the repetitions
def add_colored_google_text(paragraph, text, font_name, font_size, is_bold, default_color_rgb, use_special_colors=True):
    # Adds text to a paragraph, finding "Google" case-insensitively, coloring it, ensuring it starts with a capital 'G', and is never bold.
    google_colors = {
        'g': RGBColor(0x42, 0x85, 0xF4),  # Blue
        'o1': RGBColor(0xDB, 0x44, 0x37),  # Red
        'o2': RGBColor(0xF4, 0xB4, 0x00),  # Yellow
        'g2': RGBColor(0x42, 0x85, 0xF4),  # Blue
        'l': RGBColor(0x0F, 0x9D, 0x58),  # Green
        'e': RGBColor(0xDB, 0x44, 0x37)  # Red
    }
    parts = re.split(r'(google)', text, flags=re.IGNORECASE)
    for part in parts:
        if part.lower() == 'google':
            if use_special_colors:
                run = paragraph.add_run()
                run.text = 'G'
                run.font.color.rgb = google_colors['g']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False

                run = paragraph.add_run()
                run.text = 'o'
                run.font.color.rgb = google_colors['o1']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False

                run = paragraph.add_run()
                run.text = 'o'
                run.font.color.rgb = google_colors['o2']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False

                run = paragraph.add_run()
                run.text = 'g'
                run.font.color.rgb = google_colors['g2']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False

                run = paragraph.add_run()
                run.text = 'l'
                run.font.color.rgb = google_colors['l']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False

                run = paragraph.add_run()
                run.text = 'e'
                run.font.color.rgb = google_colors['e']
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False
            else:
                # Add "Google" with the default color, black
                run = paragraph.add_run()
                run.text = 'Google'
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = False
                run.font.color.rgb = default_color_rgb
        else:
            # The texts that are not "Google"
            if part:
                run = paragraph.add_run()
                run.text = part
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.color.rgb = default_color_rgb

def add_background_image_with_brightness(slide, image_path, brightness_factor=0.9):
    # Opens an image, adjusts its brightness, and adds it as a background
    try:
        image = Image.open(image_path)
        enhancer = ImageEnhance.Brightness(image)
        enhanced_image = enhancer.enhance(brightness_factor)
        image_stream = io.BytesIO()
        enhanced_image.save(image_stream, format='PNG')
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
    except FileNotFoundError:
        print(f"Error: Background image not found at {image_path}")


def add_cropped_picture(slide, image_path, left, top, crop_right_percent=0.4, **kwargs):
    # Opens an image, crops a percentage from the right side, and adds it to the slide.
    try:
        image = Image.open(image_path)
        original_width, original_height = image.size
        new_right = original_width * (1 - crop_right_percent)
        crop_box = (0, 0, new_right, original_height)
        cropped_image = image.crop(crop_box)
        image_stream = io.BytesIO()
        cropped_image.save(image_stream, format='PNG')
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, **kwargs)

    except FileNotFoundError:
        print(f"Error: Image not found at {image_path}")


# Presentation Setup
prs = pptx.Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6]  # Blank layout

# Font settings and colors
BLACK_FONT = RGBColor(0x20, 0x21, 0x24)
WHITE_FONT = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Open Sans"

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide1, 'title_slide.png')

# Add Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
p_title = title_box.text_frame.paragraphs[0]
p_title.alignment = PP_ALIGN.CENTER
add_colored_google_text(p_title, "Why Google Glass Failed", FONT_NAME, 44, True, BLACK_FONT, use_special_colors=True)

# Add Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), Inches(1))
p_subtitle = subtitle_box.text_frame.paragraphs[0]
p_subtitle.alignment = PP_ALIGN.CENTER
p_subtitle.text = "An analysis of the ambitious project's rise and fall"
p_subtitle.font.name = FONT_NAME
p_subtitle.font.size = Pt(24)
p_subtitle.font.color.rgb = BLACK_FONT

# Add Navigation Circles. At first I was going to add transitions but that is not available in this library
circle_topics = [
    ("What Was Google Glass?", RGBColor(0x4A, 0x90, 0xE2)),  # Blue
    ("The Hype vs. The Reality", RGBColor(0x4A, 0x90, 0xE2)),  # Blue
    ("The Price Point", RGBColor(0xD0, 0x02, 0x1B)),  # Red
    ("Major Privacy Concerns", RGBColor(0xD0, 0x02, 0x1B)),  # Red
    ("Awkward Design", RGBColor(0xF5, 0xA6, 0x23)),  # Yellow
    ("Lack of a Clear Purpose", RGBColor(0xF5, 0xA6, 0x23)),  # Yellow
    ("The Glasshole Effect", RGBColor(0x7E, 0xD3, 0x21)),  # Green
    ("The Enterprise Edition: A Second Life?", RGBColor(0x7E, 0xD3, 0x21))  # Green
]

# Positions for the circles
circle_size = 2.0
gap_size = 1.8  # gap size for horizontal spread
total_width = circle_size * 4 + gap_size * 3
start_x = (16 - total_width) / 2
positions = [
    (start_x, 4.5), (start_x + circle_size + gap_size, 4.5), (start_x + 2 * (circle_size + gap_size), 4.5),
    (start_x + 3 * (circle_size + gap_size), 4.5),
    (start_x, 6.7), (start_x + circle_size + gap_size, 6.7), (start_x + 2 * (circle_size + gap_size), 6.7),
    (start_x + 3 * (circle_size + gap_size), 6.7)
]

for i, (topic, color) in enumerate(circle_topics):
    left, top = positions[i]
    shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(circle_size), Inches(circle_size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE_FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

# Slide 2: What Was Google Glass?
slide2 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide2, 'blue_bg.png') # Background color for slide

# Title
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1))
p2_title = title2.text_frame.paragraphs[0]
p2_title.alignment = PP_ALIGN.LEFT
add_colored_google_text(p2_title, "What Was Google Glass?", FONT_NAME, 36, True, WHITE_FONT, use_special_colors=False)

# Content
content2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(7))
tf2 = content2.text_frame
tf2.word_wrap = True
tf2.clear()

text_lines_2 = [
    "A wearable, voice-controlled Android device.",
    "It featured a small screen, a camera, and a bone-conduction speaker.",
    "It was designed to provide a hands-free, augmented reality experience.",
    "The \"Explorer Edition\" was released to developers and early adopters in 2013 for $1,500."
]
for line in text_lines_2:
    p = tf2.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE_FONT
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)
    text_run = p.add_run()
    text_run.text = line
    text_run.font.size = Pt(24)

# images
slide2.shapes.add_picture('slide2_components.png', Inches(8), Inches(0), width=Inches(8),
                          height=Inches(4.5))  # Top right
slide2.shapes.add_picture('slide2_diagram.png', Inches(8), Inches(4.5), width=Inches(8),
                          height=Inches(4.5))  # Bottom right

# Slide 3: The Hype vs. The Reality
slide3 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide3, 'blue_bg.png')

# Title
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1))
p3_title = title3.text_frame.paragraphs[0]
p3_title.text = "The Hype vs. The Reality"
p3_title.font.name = FONT_NAME
p3_title.font.size = Pt(36)
p3_title.font.bold = True
p3_title.font.color.rgb = WHITE_FONT
p3_title.alignment = PP_ALIGN.LEFT

# Content
content3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(7))
tf3 = content3.text_frame
tf3.word_wrap = True
tf3.clear()  # Remove default paragraph

p3_1 = tf3.add_paragraph()
p3_1.text = "The Hype:"
p3_1.font.bold = True
p3_1.font.name = FONT_NAME
p3_1.font.size = Pt(24)
p3_1.font.color.rgb = WHITE_FONT

hype_lines = [
    "A revolutionary device that would change the way we interact with technology.",
    "Featured in high-fashion magazines and worn by celebrities.",
    "Promised a future of seamless augmented reality."
]
for line in hype_lines:
    p = tf3.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE_FONT
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)
    text_run = p.add_run()
    text_run.text = line
    text_run.font.size = Pt(24)

p3_2 = tf3.add_paragraph()
p3_2.text = "\n"
run_reality = p3_2.add_run()
run_reality.text = "The Reality:"
run_reality.font.bold = True
run_reality.font.name = FONT_NAME
run_reality.font.size = Pt(24)
run_reality.font.color.rgb = WHITE_FONT

reality_lines = [
    "A clunky, unfinished prototype.",
    "Limited functionality and poor battery life.",
    "A host of technical and social problems."
]
for line in reality_lines:
    p = tf3.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE_FONT
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)
    text_run = p.add_run()
    text_run.text = line
    text_run.font.size = Pt(24)

# Images at right corners
slide3.shapes.add_picture('slide3_fashion.png', Inches(8), Inches(0), width=Inches(8), height=Inches(4.5))
slide3.shapes.add_picture('slide3_social.jpg', Inches(8), Inches(4.5), width=Inches(8), height=Inches(4.5))

# Slide 4: The Prohibitive Price Point
slide4 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide4, 'red_bg.png')
add_cropped_picture(slide4, 'slide4_price.png', left=Inches(0), top=Inches(0), crop_right_percent=0.4, height=prs.slide_height)

title4 = slide4.shapes.add_textbox(Inches(8), Inches(0.5), Inches(7.5), Inches(1))
p4_title = title4.text_frame.paragraphs[0]
p4_title.text = "The Prohibitive Price Point"
p4_title.font.name = FONT_NAME
p4_title.font.size = Pt(36)
p4_title.font.bold = True
p4_title.font.color.rgb = WHITE_FONT
content4 = slide4.shapes.add_textbox(Inches(8), Inches(1.5), Inches(7.5), Inches(6))
tf4 = content4.text_frame
tf4.word_wrap = True
tf4.clear()

text_lines_4 = [
    "The Explorer Edition cost $1,500.",
    "This was far too expensive for the average consumer, especially for a first-generation device with limited functionality.",
    "The high price created an image of elitism and exclusivity, which alienated many potential users.",
    "The bill of materials was estimated to be around $80, which made the high price tag even more difficult to justify."
]
for line in text_lines_4:
    p = tf4.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE_FONT
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)
    text_run = p.add_run()
    text_run.text = line
    text_run.font.size = Pt(24)

# Slide 5: Sales Graph
slide4_5 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide4_5, 'red_bg.png')

title4_5 = slide4_5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
p4_5_title = title4_5.text_frame.paragraphs[0]
p4_5_title.text = "Estimated Sales Figures"
p4_5_title.font.name = FONT_NAME
p4_5_title.font.size = Pt(36)
p4_5_title.font.bold = True
p4_5_title.font.color.rgb = WHITE_FONT
p4_5_title.alignment = PP_ALIGN.CENTER

# Chart Data, from public record
chart_data = CategoryChartData()
chart_data.categories = ['Q2 2013', 'Q3 2013', 'Q4 2013', 'Q1 2014', 'Q2 2014', 'Q3 2014', 'Q4 2014']
chart_data.add_series('Estimated Units Sold', (1000, 2500, 4000, 8000, 5000, 2000, 500))

# Add and Format Chart
x, y, cx, cy = Inches(2), Inches(1.5), Inches(12), Inches(6)
graphic_frame = slide4_5.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)
chart = graphic_frame.chart

chart.has_legend = False
chart.chart_title.text_frame.text = 'Estimated Google Glass Explorer Units Sold'
chart.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE_FONT
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)

# Format Category X-axis
category_axis = chart.category_axis
category_axis.has_title = True
category_axis.axis_title.text_frame.text = 'Time Period (Quarters)'
category_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = WHITE_FONT
category_axis.tick_labels.font.color.rgb = WHITE_FONT
category_axis.format.line.fill.solid()
category_axis.format.line.fill.fore_color.rgb = WHITE_FONT

# Format Value Y-axis
value_axis = chart.value_axis
value_axis.has_title = True
value_axis.axis_title.text_frame.text = 'Estimated Units Sold'
value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = WHITE_FONT
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.fill.solid()
value_axis.major_gridlines.format.line.fill.fore_color.rgb = WHITE_FONT
value_axis.tick_labels.font.color.rgb = WHITE_FONT
value_axis.format.line.fill.solid()
value_axis.format.line.fill.fore_color.rgb = WHITE_FONT

# Format Data Labels and Colors
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(12)
data_labels.font.color.rgb = WHITE_FONT
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = RGBColor(0x42, 0x85, 0xF4)

# Slide 6: Major Privacy Concerns
slide5 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide5, 'red_bg.png')
slide5.shapes.add_picture('slide5_concerns.png', 0, 0, width=Inches(7.8), height=prs.slide_height)
title5 = slide5.shapes.add_textbox(Inches(8), Inches(0.5), Inches(7.5), Inches(1))
p5_title = title5.text_frame.paragraphs[0]
p5_title.text = "Major Privacy Concerns"
p5_title.font.name = FONT_NAME
p5_title.font.size = Pt(36)
p5_title.font.bold = True
p5_title.font.color.rgb = WHITE_FONT
content5 = slide5.shapes.add_textbox(Inches(8), Inches(1.5), Inches(7.5), Inches(6.5))
tf5 = content5.text_frame
tf5.word_wrap = True
tf5.clear()

text_lines_5 = [
    "The built-in camera raised serious privacy questions.",
    "The ability to record video and take pictures discreetly made people uncomfortable.",
    "The term \"Glasshole\" was coined to describe users who were perceived as invading the privacy of others.",
    "Many businesses, such as bars and movie theaters, banned the use of Google Glass on their premises."
]
for line in text_lines_5:
    p = tf5.add_paragraph()

    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.name = FONT_NAME
    bullet_run.font.size = Pt(28)
    bullet_run.font.color.rgb = WHITE_FONT

    add_colored_google_text(p, line, FONT_NAME, 24, False, WHITE_FONT, use_special_colors=False)

# Slide 7: Awkward and Unfashionable Design
slide6 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide6, 'yellow_bg.png')
title6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
p6_title = title6.text_frame.paragraphs[0]
p6_title.text = "Awkward and Unfashionable Design"
p6_title.font.name = FONT_NAME
p6_title.font.size = Pt(36)
p6_title.font.bold = True
p6_title.font.color.rgb = WHITE_FONT
p6_title.alignment = PP_ALIGN.CENTER
content6 = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(3))
tf6 = content6.text_frame
tf6.word_wrap = True
tf6.clear()
text_lines_6 = [
    "The device was bulky, lopsided, and generally considered to be unattractive.",
    "It was not a device that people felt comfortable wearing in public.",
    "The design screamed \"tech gadget\" rather than \"fashion accessory.\"",
    "Despite collaborations with fashion designers, the fundamental design remained a major turn-off for consumers."
]
for line in text_lines_6:
    p = tf6.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE_FONT
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)
    text_run = p.add_run()
    text_run.text = line
    text_run.font.size = Pt(24)

slide6.shapes.add_picture('slide6_awkward.jpg', 0, Inches(4.5), width=prs.slide_width)

# Slide 8: Lack of a Clear Purpose
slide7 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide7, 'yellow_bg.png')
title7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
p7_title = title7.text_frame.paragraphs[0]
p7_title.text = "Lack of a Clear Purpose"
p7_title.font.name = FONT_NAME
p7_title.font.size = Pt(36)
p7_title.font.bold = True
p7_title.font.color.rgb = WHITE_FONT
p7_title.alignment = PP_ALIGN.CENTER
content7 = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(3.5))
tf7 = content7.text_frame
tf7.word_wrap = True
tf7.clear()

text_lines_7 = [
    "Google Glass was a solution in search of a problem.",
    "There was no single, compelling use case that made the device a \"must-have.\"",
    "The functionality it offered could be easily replicated by a smartphone.",
    "Without a clear purpose, it was difficult for consumers to justify the high price and social awkwardness."
]
for line in text_lines_7:
    p = tf7.add_paragraph()
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.name = FONT_NAME
    bullet_run.font.size = Pt(28)
    bullet_run.font.color.rgb = WHITE_FONT
    add_colored_google_text(p, line, FONT_NAME, 24, False, WHITE_FONT, use_special_colors=False)

# Center image, make it large
img_height_7 = Inches(4.0)
pic7 = slide7.shapes.add_picture('question_glass.png', Inches(0), Inches(4.65), height=img_height_7)
pic7.left = int((prs.slide_width - pic7.width) / 2)

# Slide 9: Social Awkwardness
slide8 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide8, 'green_bg.png')
title8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
p8_title = title8.text_frame.paragraphs[0]
p8_title.alignment = PP_ALIGN.CENTER
add_colored_google_text(p8_title, "Social Awkwardness and the \"Glasshole\" Effect", FONT_NAME, 36, True, WHITE_FONT, use_special_colors=False)

content8 = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(3.5))
tf8 = content8.text_frame
tf8.word_wrap = True
tf8.clear()

text_lines_8 = [
    "Wearing Google Glass in public was a socially awkward experience.",
    "It created a barrier between the user and the people they were interacting with.",
    "The device was seen as a sign of social disengagement and a potential invasion of privacy.",
    "The negative social stigma was a major deterrent for potential users."
]

for line in text_lines_8:
    p = tf8.add_paragraph()
    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.name = FONT_NAME
    bullet_run.font.size = Pt(28)
    bullet_run.font.color.rgb = WHITE_FONT
    add_colored_google_text(p, line, FONT_NAME, 24, False, WHITE_FONT, use_special_colors=False)

slide8.shapes.add_picture('slide8_glasshole.jpg', Inches(4), Inches(4.7), width=Inches(8))

# Slide 10: The Enterprise Edition
slide9 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide9, 'green_bg.png')

# Title
title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1.5))
p9_title = title9.text_frame.paragraphs[0]
p9_title.alignment = PP_ALIGN.LEFT
add_colored_google_text(p9_title, "The Enterprise Edition:\nA Second Life?", FONT_NAME, 36, True, WHITE_FONT, use_special_colors=False)

# Content
content9 = slide9.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7.5), Inches(6.5))
tf9 = content9.text_frame
tf9.word_wrap = True
tf9.clear()

text_lines_9 = [
    "Google eventually discontinued the consumer version of Glass in 2015.",
    "However, the company pivoted to the enterprise market with the \"Glass Enterprise Edition.\"",
    "This version of the device has found success in industries like manufacturing, logistics, and healthcare.",
    "In these contexts, the hands-free functionality and augmented reality features provide real value.",
    "But even those efforts were ultimately discontinued due to the tool's low impact on daily life."
]
for line in text_lines_9:
    p = tf9.add_paragraph()

    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.name = FONT_NAME
    bullet_run.font.size = Pt(28)
    bullet_run.font.color.rgb = WHITE_FONT

    add_colored_google_text(p, line, FONT_NAME, 24, False, WHITE_FONT, use_special_colors=False)

# Images at right corners
slide9.shapes.add_picture('slide9_edition.jpg', Inches(8), Inches(0), width=Inches(8), height=Inches(4.5))
slide9.shapes.add_picture('slide9_person.png', Inches(8), Inches(4.5), width=Inches(8), height=Inches(4.5))

# Slide 11: Lessons Learned
slide10 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide10, 'conc_slide.png')
title10 = slide10.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
p10_title = title10.text_frame.paragraphs[0]
p10_title.text = "Lessons Learned"
p10_title.font.name = FONT_NAME
p10_title.font.size = Pt(44)
p10_title.font.bold = True
p10_title.font.color.rgb = BLACK_FONT
p10_title.alignment = PP_ALIGN.CENTER
content10 = slide10.shapes.add_textbox(Inches(1.5), Inches(1.75), Inches(13), Inches(5))
tf10 = content10.text_frame
tf10.word_wrap = True
tf10.clear()

text_lines_10 = [
    ("Don't release a prototype as a consumer product:",
     " The Explorer Edition was not ready for the public, and the negative first impression was difficult to overcome."),
    ("Privacy is paramount:",
     " In the age of connected devices, privacy must be a primary consideration in product design."),
    ("Design and social acceptability matter:",
     " For wearable technology, fashion and social norms are just as important as functionality."),
    ("A clear value proposition is essential:",
     " A new product needs to solve a real problem or offer a compelling new experience to succeed.")
]
for bold_part, regular_part in text_lines_10:
    p = tf10.add_paragraph()
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK_FONT

    bullet_run = p.add_run()
    bullet_run.text = "∙ "
    bullet_run.font.size = Pt(28)

    bold_run = p.add_run()
    bold_run.text = bold_part
    bold_run.font.bold = True
    bold_run.font.size = Pt(24)

    regular_run = p.add_run()
    regular_run.text = regular_part
    regular_run.font.size = Pt(24)

# Google logo at the bottom center. Also scale and dimensions
logo_height = Inches(1.5)
pic = slide10.shapes.add_picture('logo.png', Inches(0), Inches(0), height=logo_height)
pic.left = int((prs.slide_width - pic.width) / 2)
pic.top = Inches(6.8)

# Slide 12: Thank You
slide11 = prs.slides.add_slide(blank_slide_layout)
add_background_image_with_brightness(slide11, 'conc_slide.png')

# text
thank_you_box = slide11.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
p_thank_you = thank_you_box.text_frame.paragraphs[0]
p_thank_you.text = "Thank You"
p_thank_you.font.name = FONT_NAME
p_thank_you.font.size = Pt(66)
p_thank_you.font.bold = True
p_thank_you.font.color.rgb = BLACK_FONT
p_thank_you.alignment = PP_ALIGN.CENTER

# Add by: names...
names_box = slide11.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
p_names = names_box.text_frame.paragraphs[0]
p_names.text = "Made by: Ahmed Ismail, Manil Laroussi, and Ahmed Al Ali"
p_names.font.name = FONT_NAME
p_names.font.size = Pt(24)
p_names.font.bold = True
p_names.font.color.rgb = BLACK_FONT
p_names.alignment = PP_ALIGN.CENTER

# Save Presentation
try:
    prs.save("Google_Glass_Failure_Presentation.pptx")
    print("Presentation 'Google_Glass_Failure_Presentation.pptx' created successfully.")
except Exception as e:
    print(f"An error occurred while saving the presentation: {e}")

